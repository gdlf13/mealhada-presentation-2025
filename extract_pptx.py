import os
import sys
import json
from pptx import Presentation
from pathlib import Path
import re
import shutil

def clean_text(text):
    """Clean text from PowerPoint, removing extra spaces and newlines."""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def extract_pptx_content(pptx_path, output_dir):
    """Extract content from PowerPoint and save as structured data."""
    prs = Presentation(pptx_path)
    
    # Create output directories
    images_dir = Path(output_dir) / "public" / "images" / "slides"
    content_dir = Path(output_dir) / "src" / "content" / "modules"
    
    images_dir.mkdir(parents=True, exist_ok=True)
    content_dir.mkdir(parents=True, exist_ok=True)
    
    # Structure to hold all slide data
    presentation_data = {
        "title": "14 Encontro com a Educação Mealhada", # Updated title
        "modules": [],
        "totalSlides": len(prs.slides)
    }
    
    current_module = None
    modules = []
    module_index = 0
    
    # Process each slide
    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        slide_data = {
            "id": slide_number,
            "title": "",
            "content": [],
            "hasImage": False,
            "images": []
        }
        
        # Extract title if available
        if slide.shapes.title:
            title_text = clean_text(slide.shapes.title.text)
            slide_data["title"] = title_text
            
            # Check if this is a module title slide
            if slide_number == 1 or (title_text and any(kw in title_text.lower() for kw in ['módulo', 'secção', 'parte', 'tema'])):
                if current_module:
                    modules.append(current_module)
                
                module_index += 1
                module_id = f"module-{module_index}"
                current_module = {
                    "id": module_id,
                    "title": title_text,
                    "slides": [],
                    "description": ""
                }
        
        # Process all shapes (text boxes, images, etc.)
        for shape in slide.shapes:
            # Process text
            if hasattr(shape, "text") and shape.text:
                text = clean_text(shape.text)
                if text and text != slide_data["title"]:  # Don't duplicate the title
                    slide_data["content"].append({"type": "text", "value": text})
                    
                    # If this is the first content slide of a module, use it as description
                    if current_module and not current_module["description"] and len(current_module["slides"]) <= 1:
                        current_module["description"] = text[:200] + ("..." if len(text) > 200 else "")
            
            # Process images
            if shape.shape_type == 13:  # 13 is the enum value for pictures
                slide_data["hasImage"] = True
                image_filename = f"slide_{slide_number}_image_{len(slide_data['images']) + 1}.png"
                image_path = images_dir / image_filename
                
                # Save the image
                try:
                    with open(image_path, "wb") as img_file:
                        img_file.write(shape.image.blob)
                    slide_data["images"].append(f"/images/slides/{image_filename}")
                except Exception as e:
                    print(f"Error saving image from slide {slide_number}: {e}")
        
        if current_module:
            current_module["slides"].append(slide_data)
    
    # Add the last module
    if current_module:
        modules.append(current_module)
    
    presentation_data["modules"] = modules
    
    # Save the extracted data as JSON
    with open(content_dir / "presentation_data.json", "w", encoding="utf-8") as f:
        json.dump(presentation_data, f, ensure_ascii=False, indent=2)
    
    # Create individual module files
    for module in modules:
        module_file = content_dir / f"{module['id']}.json"
        with open(module_file, "w", encoding="utf-8") as f:
            json.dump(module, f, ensure_ascii=False, indent=2)
    
    return presentation_data

if __name__ == "__main__":
    # Default to the specific presentation name for this project
    default_pptx = "14 Encontro com a Educação Mealhada Abril 2025.pptx"
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else default_pptx
    output_dir = "."  # Current directory
    
    if not Path(pptx_path).exists():
        print(f"Error: Presentation file '{pptx_path}' not found in the current directory.")
        sys.exit(1)
        
    print(f"Extracting content from {pptx_path}...")
    data = extract_pptx_content(pptx_path, output_dir)
    
    print(f"Extracted {len(data['modules'])} modules with {data['totalSlides']} slides total.")
    print("Content saved to src/content/modules/")
    print("Images saved to public/images/slides/")
